"""Report Builder - Generates PPTX from analysis and narrative content."""

import os
from datetime import datetime
from uuid import uuid4
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ============ Configuration ============

OUTPUT_DIR = "outputs"
os.makedirs(OUTPUT_DIR, exist_ok=True)


# ============ Template Themes ============

TEMPLATES = {
    "corporate_blue": {
        "name": "Corporate Blue",
        "primary": RGBColor(0x00, 0x52, 0xcc),
        "secondary": RGBColor(0x00, 0x96, 0xff),
        "accent": RGBColor(0x10, 0xb9, 0x81),
        "dark": RGBColor(0x1e, 0x29, 0x3b),
        "light": RGBColor(0xf0, 0xf4, 0xf8),
        "white": RGBColor(0xff, 0xff, 0xff),
        "warning": RGBColor(0xf5, 0x9e, 0x0b),
    },
    "modern_dark": {
        "name": "Modern Dark",
        "primary": RGBColor(0x63, 0x66, 0xf1),
        "secondary": RGBColor(0x8b, 0x5c, 0xf6),
        "accent": RGBColor(0x22, 0xd3, 0xee),
        "dark": RGBColor(0x0f, 0x17, 0x2a),
        "light": RGBColor(0x1e, 0x29, 0x3b),
        "white": RGBColor(0xf8, 0xfa, 0xfc),
        "warning": RGBColor(0xfb, 0xbf, 0x24),
    },
    "fresh_green": {
        "name": "Fresh Green",
        "primary": RGBColor(0x05, 0x96, 0x69),
        "secondary": RGBColor(0x10, 0xb9, 0x81),
        "accent": RGBColor(0x06, 0xb6, 0xd4),
        "dark": RGBColor(0x1a, 0x2e, 0x2a),
        "light": RGBColor(0xec, 0xfd, 0xf5),
        "white": RGBColor(0xff, 0xff, 0xff),
        "warning": RGBColor(0xf5, 0x9e, 0x0b),
    },
    "warm_orange": {
        "name": "Warm Orange",
        "primary": RGBColor(0xea, 0x58, 0x0c),
        "secondary": RGBColor(0xf9, 0x73, 0x16),
        "accent": RGBColor(0x06, 0xb6, 0xd4),
        "dark": RGBColor(0x27, 0x1c, 0x1c),
        "light": RGBColor(0xff, 0xf7, 0xed),
        "white": RGBColor(0xff, 0xff, 0xff),
        "warning": RGBColor(0xfb, 0xbf, 0x24),
    },
    "minimal_gray": {
        "name": "Minimal Gray",
        "primary": RGBColor(0x37, 0x41, 0x51),
        "secondary": RGBColor(0x6b, 0x72, 0x80),
        "accent": RGBColor(0x3b, 0x82, 0xf6),
        "dark": RGBColor(0x11, 0x18, 0x27),
        "light": RGBColor(0xf9, 0xfa, 0xfb),
        "white": RGBColor(0xff, 0xff, 0xff),
        "warning": RGBColor(0xf5, 0x9e, 0x0b),
    },
}


# ============ Presentation Builder ============

class ReportGenerator:
    """Generates PPTX presentations with rich content."""
    
    def __init__(self, template_name: str = "corporate_blue"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.template = TEMPLATES.get(template_name, TEMPLATES["corporate_blue"])
        self.template_name = template_name
        self.slides_created = []
    
    def _add_title_slide(self, title: str, subtitle: str):
        """Create polished title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Dark background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.template["dark"]
        bg.line.fill.background()
        
        # Top accent bar
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.15))
        top.fill.solid()
        top.fill.fore_color.rgb = self.template["accent"]
        top.line.fill.background()
        
        # Left accent
        left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8), Inches(0.4), Inches(1.8))
        left.fill.solid()
        left.fill.fore_color.rgb = self.template["primary"]
        left.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.template["white"]
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = self.template["secondary"]
        
        # Bottom bar
        bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), Inches(13.333), Inches(0.15))
        bottom.fill.solid()
        bottom.fill.fore_color.rgb = self.template["secondary"]
        bottom.line.fill.background()
        
        self.slides_created.append("Title")
    
    def _add_content_slide(self, title: str, content: str, bullets: list = None):
        """Create content slide with text and optional bullets."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.4))
        header.fill.solid()
        header.fill.fore_color.rgb = self.template["primary"]
        header.line.fill.background()
        
        # Accent line
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.4), Inches(13.333), Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.template["accent"]
        accent.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.template["white"]
        
        # Content area
        if content and len(content) > 50:
            # Long text - show as paragraph
            content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12), Inches(2))
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            # Truncate if too long
            display_text = content[:600] + "..." if len(content) > 600 else content
            p.text = display_text
            p.font.size = Pt(16)
            p.font.color.rgb = self.template["dark"]
            start_y = 4.0
        else:
            start_y = 1.8
        
        # Bullets
        if bullets:
            bullet_box = slide.shapes.add_textbox(Inches(0.6), Inches(start_y), Inches(12), Inches(7.5 - start_y - 0.5))
            tf = bullet_box.text_frame
            tf.word_wrap = True
            
            for i, bullet in enumerate(bullets[:6]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"→  {bullet}"
                p.font.size = Pt(18)
                p.font.color.rgb = self.template["dark"]
                p.space_before = Pt(12)
                p.space_after = Pt(8)
        
        self.slides_created.append(title)
    
    def _add_metrics_slide(self, title: str, metrics: list):
        """Create metrics slide with cards."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.4))
        header.fill.solid()
        header.fill.fore_color.rgb = self.template["secondary"]
        header.line.fill.background()
        
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.4), Inches(13.333), Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.template["accent"]
        accent.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.template["white"]
        
        # Metrics cards
        metrics = metrics[:4] if metrics else []
        if not metrics:
            metrics = [{"name": "No metrics available", "value": "N/A"}]
        
        card_width = 2.9
        gap = 0.3
        total = len(metrics) * card_width + (len(metrics) - 1) * gap
        start_x = (13.333 - total) / 2
        
        for i, m in enumerate(metrics):
            x = start_x + i * (card_width + gap)
            
            # Card
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(card_width), Inches(2.5))
            card.fill.solid()
            card.fill.fore_color.rgb = self.template["light"]
            card.line.color.rgb = self.template["secondary"]
            card.line.width = Pt(2)
            
            # Top accent
            top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.2), Inches(card_width), Inches(0.12))
            top.fill.solid()
            top.fill.fore_color.rgb = self.template["primary"]
            top.line.fill.background()
            
            # Value
            val_box = slide.shapes.add_textbox(Inches(x), Inches(2.6), Inches(card_width), Inches(1))
            tf = val_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(m.get("value", "N/A"))
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.template["primary"]
            p.alignment = PP_ALIGN.CENTER
            
            # Name
            name_box = slide.shapes.add_textbox(Inches(x), Inches(3.7), Inches(card_width), Inches(0.8))
            tf = name_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(m.get("name", "Metric"))
            p.font.size = Pt(14)
            p.font.color.rgb = self.template["dark"]
            p.alignment = PP_ALIGN.CENTER
        
        self.slides_created.append(title)
    
    def _add_comparison_slide(self, title: str, left_title: str, left_items: list, 
                               right_title: str, right_items: list):
        """Create two-column comparison slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.4))
        header.fill.solid()
        header.fill.fore_color.rgb = self.template["primary"]
        header.line.fill.background()
        
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.4), Inches(13.333), Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.template["accent"]
        accent.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.template["white"]
        
        # Left card
        left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.7), Inches(6.1), Inches(5.4))
        left_card.fill.solid()
        left_card.fill.fore_color.rgb = self.template["light"]
        left_card.line.fill.background()
        
        left_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.7), Inches(6.1), Inches(0.6))
        left_header.fill.solid()
        left_header.fill.fore_color.rgb = self.template["secondary"]
        left_header.line.fill.background()
        
        lh_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(0.5))
        tf = lh_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"✓ {left_title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.template["white"]
        
        left_content = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(5.8), Inches(4.4))
        tf = left_content.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items[:5] if left_items else ["No data"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(15)
            p.font.color.rgb = self.template["dark"]
            p.space_before = Pt(10)
        
        # Right card
        right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(6.1), Inches(5.4))
        right_card.fill.solid()
        right_card.fill.fore_color.rgb = self.template["light"]
        right_card.line.fill.background()
        
        right_header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.7), Inches(6.1), Inches(0.6))
        right_header.fill.solid()
        right_header.fill.fore_color.rgb = self.template["warning"]
        right_header.line.fill.background()
        
        rh_box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5.8), Inches(0.5))
        tf = rh_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"⚠ {right_title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.template["white"]
        
        right_content = slide.shapes.add_textbox(Inches(7), Inches(2.5), Inches(5.8), Inches(4.4))
        tf = right_content.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items[:5] if right_items else ["No data"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(15)
            p.font.color.rgb = self.template["dark"]
            p.space_before = Pt(10)
        
        self.slides_created.append(title)
    
    def _add_recommendations_slide(self, title: str, recommendations: list):
        """Create recommendations slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.4))
        header.fill.solid()
        header.fill.fore_color.rgb = self.template["dark"]
        header.line.fill.background()
        
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.4), Inches(13.333), Inches(0.08))
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.template["accent"]
        accent.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.template["white"]
        
        recs = recommendations[:5] if recommendations else ["No recommendations available"]
        
        for i, rec in enumerate(recs):
            y = 1.8 + i * 1.05
            
            # Number circle
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y + 0.1), Inches(0.5), Inches(0.5))
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.template["primary"]
            circle.line.fill.background()
            
            num_box = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.15), Inches(0.5), Inches(0.4))
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.template["white"]
            p.alignment = PP_ALIGN.CENTER
            
            # Text
            rec_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11.5), Inches(0.9))
            tf = rec_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = rec
            p.font.size = Pt(17)
            p.font.color.rgb = self.template["dark"]
        
        self.slides_created.append(title)
    
    def save(self, filepath: str) -> str:
        self.prs.save(filepath)
        return filepath


# ============ Content Extraction Helpers ============

def extract_summary_bullets(narrative: dict) -> list:
    """Extract bullet points from executive summary."""
    summary = narrative.get("executive_summary", "")
    if not summary:
        return ["Analysis complete", "See detailed metrics below"]
    
    # Split into sentences and take top 4
    sentences = [s.strip() for s in summary.replace("\n", " ").split(".") if len(s.strip()) > 20]
    return sentences[:4] if sentences else [summary[:200]]


def extract_metrics(analysis: dict) -> list:
    """Extract key metrics for display."""
    metrics = []
    
    # From summary metrics
    for m in analysis.get("summary_metrics", [])[:2]:
        metrics.append({"name": m.get("name", "Metric"), "value": m.get("value", "N/A")})
    
    # From performance metrics
    for m in analysis.get("performance_metrics", [])[:2]:
        metrics.append({"name": m.get("name", "Metric"), "value": m.get("value", "N/A")})
    
    if not metrics:
        # Fallback - generate from dataset meta
        meta = analysis.get("dataset_meta", {})
        if meta:
            metrics.append({"name": "Total Rows", "value": str(meta.get("total_rows", "N/A"))})
            metrics.append({"name": "Total Columns", "value": str(meta.get("total_columns", "N/A"))})
    
    return metrics[:4]


def extract_performers(analysis: dict) -> tuple:
    """Extract top and bottom performers."""
    top = []
    bottom = []
    
    for p in analysis.get("top_performers", []):
        name = p.get("entity_name", "Unknown")
        value = p.get("metric_value", "")
        top.append(f"{name}: {value}" if value else name)
    
    for p in analysis.get("bottom_performers", []):
        name = p.get("entity_name", "Unknown")
        value = p.get("metric_value", "")
        bottom.append(f"{name}: {value}" if value else name)
    
    # Fallbacks
    if not top:
        top = ["No top performers identified"]
    if not bottom:
        bottom = ["No underperformers identified"]
    
    return top[:5], bottom[:5]


def extract_trends_and_risks(narrative: dict, analysis: dict) -> list:
    """Extract trends and risks as bullet points."""
    bullets = []
    
    # From narrative
    trends = narrative.get("trend_commentary", "")
    if trends and len(trends) > 20:
        sentences = [s.strip() for s in trends.split(".") if len(s.strip()) > 15]
        bullets.extend(sentences[:2])
    
    risks = narrative.get("risk_and_anomalies", "")
    if risks and len(risks) > 20:
        sentences = [s.strip() for s in risks.split(".") if len(s.strip()) > 15]
        bullets.extend(sentences[:2])
    
    # From analysis
    for a in analysis.get("anomalies", [])[:2]:
        bullets.append(f"⚠ {a.get('title', 'Anomaly')}: {a.get('description', '')[:50]}")
    
    for t in analysis.get("trends", [])[:2]:
        bullets.append(f"📈 {t.get('metric', 'Metric')}: {t.get('direction', '')} ({t.get('change_value', '')})")
    
    if not bullets:
        bullets = ["No significant trends or risks detected", "Data appears stable"]
    
    return bullets[:5]


def extract_recommendations(narrative: dict) -> list:
    """Extract recommendations."""
    recs = narrative.get("recommendations", [])
    if recs and isinstance(recs, list):
        return [str(r) for r in recs[:5]]
    return ["Continue monitoring key metrics", "Review data regularly", "Investigate any anomalies"]


def auto_select_template(analysis: dict, narrative: dict) -> str:
    """Auto-select best template based on content."""
    summary = narrative.get("executive_summary", "").lower()
    
    if any(word in summary for word in ["growth", "increase", "success", "improve"]):
        return "fresh_green"
    elif any(word in summary for word in ["decline", "drop", "concern", "risk"]):
        return "warm_orange"
    elif any(word in summary for word in ["tech", "data", "analytics", "digital"]):
        return "modern_dark"
    else:
        return "corporate_blue"


# ============ Main Builder Function ============

async def build_report(analysis: dict, narrative: dict) -> dict:
    """Build a complete PPTX report from analysis and narrative.
    
    This function GUARANTEES 6 content-rich slides.
    """
    # Select template
    template = auto_select_template(analysis, narrative)
    generator = ReportGenerator(template)
    
    # Get dataset info for title
    dataset_id = analysis.get("dataset_id", "Dataset")
    meta = analysis.get("dataset_meta", {})
    rows = meta.get("total_rows", "N/A")
    cols = meta.get("total_columns", "N/A")
    
    # SLIDE 1: Title
    generator._add_title_slide(
        title="H-OOI Analysis Report",
        subtitle=f"Dataset: {dataset_id} | {rows} rows × {cols} columns | {datetime.now().strftime('%B %d, %Y')}"
    )
    
    # SLIDE 2: Executive Summary
    summary_text = narrative.get("executive_summary", "Analysis complete.")
    summary_bullets = extract_summary_bullets(narrative)
    generator._add_content_slide(
        title="Executive Summary",
        content=summary_text,
        bullets=summary_bullets
    )
    
    # SLIDE 3: Key Metrics
    metrics = extract_metrics(analysis)
    generator._add_metrics_slide(
        title="Key Metrics",
        metrics=metrics
    )
    
    # SLIDE 4: Performance Analysis
    top, bottom = extract_performers(analysis)
    generator._add_comparison_slide(
        title="Performance Analysis",
        left_title="Top Performers",
        left_items=top,
        right_title="Needs Attention",
        right_items=bottom
    )
    
    # SLIDE 5: Trends & Risks
    trends_risks = extract_trends_and_risks(narrative, analysis)
    generator._add_content_slide(
        title="Trends & Risks",
        content=None,
        bullets=trends_risks
    )
    
    # SLIDE 6: Recommendations
    recs = extract_recommendations(narrative)
    generator._add_recommendations_slide(
        title="Recommendations & Next Steps",
        recommendations=recs
    )
    
    # Save
    job_id = uuid4().hex[:8]
    output_path = os.path.join(OUTPUT_DIR, f"report_{job_id}.pptx")
    generator.save(output_path)
    
    return {
        "status": "complete",
        "output_path": output_path,
        "slides_created": generator.slides_created,
        "template": template,
        "job_id": job_id
    }
